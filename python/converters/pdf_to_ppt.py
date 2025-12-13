"""PDF to PowerPoint conversion."""

import os
from typing import Optional
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

from pdf_operations import (
    SUCCESS,
    ERROR_FILE_NOT_FOUND,
    ERROR_CONVERSION_FAILED,
)


def pdf_to_pptx(
    input_path: str, 
    output_path: str,
    dpi: int = 150
) -> int:
    """Convert PDF to PowerPoint presentation.
    
    Each PDF page becomes a slide with the page rendered as an image.
    
    Args:
        input_path: Path to the PDF file.
        output_path: Path for the output PPTX file.
        dpi: Resolution for page rendering (default 150).
        
    Returns:
        SUCCESS (0) or error code.
    """
    try:
        if not os.path.exists(input_path):
            return ERROR_FILE_NOT_FOUND
            
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
            
        doc = fitz.open(input_path)
        prs = Presentation()
        
        # Set slide dimensions (widescreen 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Use blank layout
        blank_layout = prs.slide_layouts[6]
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Render page to image
            zoom = dpi / 72
            matrix = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=matrix)
            
            # Convert to bytes
            img_data = pix.tobytes("png")
            img_stream = BytesIO(img_data)
            
            # Add slide
            slide = prs.slides.add_slide(blank_layout)
            
            # Calculate image placement to fit slide
            page_ratio = pix.width / pix.height
            slide_ratio = prs.slide_width / prs.slide_height
            
            if page_ratio > slide_ratio:
                # Width-constrained
                img_width = prs.slide_width
                img_height = int(prs.slide_width / page_ratio)
                left = 0
                top = (prs.slide_height - img_height) // 2
            else:
                # Height-constrained
                img_height = prs.slide_height
                img_width = int(prs.slide_height * page_ratio)
                left = (prs.slide_width - img_width) // 2
                top = 0
                
            slide.shapes.add_picture(img_stream, left, top, img_width, img_height)
            
        doc.close()
        prs.save(output_path)
        
        return SUCCESS
        
    except FileNotFoundError:
        return ERROR_FILE_NOT_FOUND
    except Exception:
        return ERROR_CONVERSION_FAILED
