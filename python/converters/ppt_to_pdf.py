"""PowerPoint to PDF converter using python-pptx and reportlab."""

import os
from typing import List

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from PIL import Image
    import io
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


SUCCESS = 0
ERROR_FILE_NOT_FOUND = -1
ERROR_CONVERSION_FAILED = -3
ERROR_MISSING_DEPENDENCY = -6


def ppt_to_pdf(input_path: str, output_path: str) -> int:
    """Convert PowerPoint presentation to PDF.
    
    Args:
        input_path: Path to PowerPoint file (.ppt, .pptx).
        output_path: Output PDF path.
        
    Returns:
        SUCCESS (0) or error code.
    """
    if not HAS_PPTX or not HAS_REPORTLAB or not HAS_PIL:
        return ERROR_MISSING_DEPENDENCY
    
    try:
        if not os.path.exists(input_path):
            return ERROR_FILE_NOT_FOUND
        
        # Load presentation
        prs = Presentation(input_path)
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Calculate PDF page size (landscape)
        pdf_width = 792  # 11 inches
        pdf_height = 612  # 8.5 inches
        
        # Create PDF
        c = canvas.Canvas(output_path, pagesize=(pdf_width, pdf_height))
        
        for i, slide in enumerate(prs.slides):
            # Draw slide number and title
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, pdf_height - 50, f"Slide {i + 1}")
            
            # Try to extract title
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                c.setFont("Helvetica-Bold", 24)
                c.drawString(50, pdf_height - 100, title_text[:60])
            
            # Extract text from all text frames
            y_pos = pdf_height - 150
            c.setFont("Helvetica", 12)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Wrap long text
                    lines = text.split('\n')
                    for line in lines[:15]:  # Limit lines per slide
                        if len(line) > 80:
                            line = line[:80] + "..."
                        c.drawString(50, y_pos, line)
                        y_pos -= 20
                        if y_pos < 50:
                            break
                    if y_pos < 50:
                        break
            
            # Add new page for next slide
            if i < len(prs.slides) - 1:
                c.showPage()
        
        c.save()
        return SUCCESS
        
    except FileNotFoundError:
        return ERROR_FILE_NOT_FOUND
    except Exception as e:
        print(f"Error converting PPT to PDF: {e}")
        return ERROR_CONVERSION_FAILED
